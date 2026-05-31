"""
產出機器視覺檢測期末專題簡報 (.pptx)
風格：極簡 navy + 灰階，無 emoji，大字體
輸出：outputs/PCB_AOI_期末簡報.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent
FIG = REPO / "outputs" / "figures"
OUT = REPO / "docs" / "PCB_AOI_期末簡報.pptx"

# ---------- 極簡色票 ----------
NAVY        = RGBColor(0x1E, 0x3A, 0x5F)   # 主色（標題、強調）
NAVY_DARK   = RGBColor(0x0F, 0x1F, 0x35)
DARK        = RGBColor(0x1F, 0x29, 0x37)   # 內文
GRAY        = RGBColor(0x55, 0x60, 0x6F)   # 次要文字
GRAY_LIGHT  = RGBColor(0xCB, 0xD0, 0xD6)   # 框線
BG_ALT      = RGBColor(0xF4, 0xF5, 0xF7)   # 淺底
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT      = RGBColor(0xB1, 0x2A, 0x2A)   # 唯一強調色（僅用於標題底線、必要重點）

FONT = "Microsoft JhengHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


# ---------- 工具 ----------
def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, text, left, top, width, height, *, size=20, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, left, top, width, height, *, size=18,
                color=DARK, space_before=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(space_before)
        run = p.add_run()
        run.text = f"・ {item}"
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, *, fill=BG_ALT, line=GRAY_LIGHT,
             line_w=1.0, shape=MSO_AUTO_SHAPE_TYPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_title_bar(slide, title_text):
    """頂部標題 + navy 細底線"""
    add_text(slide, title_text, Inches(0.5), Inches(0.3), Inches(12.3),
             Inches(0.85), size=36, bold=True, color=NAVY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                  Inches(0.5), Inches(1.15),
                                  Inches(12.3), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = NAVY
    line.line.fill.background()
    line.shadow.inherit = False


def add_footer(slide, page_no, total=11):
    add_text(slide, f"{page_no} / {total}",
             Inches(0), Inches(7.1), SW, Inches(0.3),
             size=11, color=GRAY, align=PP_ALIGN.CENTER)


def add_table(slide, data, left, top, width, height, *,
              header_size=16, body_size=14, col_widths=None):
    """data: list[list[str]]，第一列為表頭"""
    rows = len(data); cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.1)
            cell.margin_top = cell.margin_bottom = Inches(0.06)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            run.font.name = FONT
            run.font.size = Pt(header_size if r == 0 else body_size)
            if r == 0:
                run.font.bold = True
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                run.font.color.rgb = DARK
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else BG_ALT
    return table


# ============================================================
# Slide 1 — 封面（純白底、navy 主視覺）
# ============================================================
def slide_1_cover():
    s = blank_slide()

    # 左側 navy 色塊
    bar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                             0, 0, Inches(0.5), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background(); bar.shadow.inherit = False

    add_text(s, "機器視覺檢測 ・ 期末專題",
             Inches(1.2), Inches(1.1), Inches(11.0), Inches(0.5),
             size=18, color=GRAY)

    add_text(s, "基於 YOLOv12 深度學習之",
             Inches(1.2), Inches(2.0), Inches(11.0), Inches(0.9),
             size=36, color=DARK)
    add_text(s, "PCB 四類缺陷自動光學檢測系統",
             Inches(1.2), Inches(2.9), Inches(11.0), Inches(1.2),
             size=48, bold=True, color=NAVY)

    add_text(s, "PCB Defect Detection via YOLOv12",
             Inches(1.2), Inches(4.3), Inches(11.0), Inches(0.6),
             size=20, color=GRAY)

    # navy 分隔細線
    line = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                              Inches(1.2), Inches(5.3),
                              Inches(2.5), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = NAVY
    line.line.fill.background(); line.shadow.inherit = False

    add_text(s, "專題成員", Inches(1.2), Inches(5.5), Inches(11.0),
             Inches(0.45), size=16, color=GRAY)
    add_text(s, "411287004   許豐有",
             Inches(1.2), Inches(6.0), Inches(11.0), Inches(0.5),
             size=22, bold=True, color=DARK)
    add_text(s, "411287035   趙梓辰",
             Inches(1.2), Inches(6.5), Inches(11.0), Inches(0.5),
             size=22, bold=True, color=DARK)


# ============================================================
# Slide 2 — 檢測物件
# ============================================================
def slide_2_object():
    s = blank_slide()
    add_title_bar(s, "檢測物件")

    add_text(s, "PCB 基板與四類常見缺陷",
             Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             size=20, color=GRAY)

    # 左：表格
    data = [
        ["缺陷類別", "形狀", "顏色", "紋理"],
        ["缺孔 Missing Hole", "應為圓孔處呈實心", "銀 / 銅", "平滑無洞"],
        ["鼠咬 Mouse Bite",   "邊緣不規則凹陷",   "銅色",    "鋸齒邊"],
        ["斷路 Open Circuit", "線路中斷成兩段",   "銅 + 綠", "高對比邊界"],
        ["短路 Short Circuit", "線路間多出橋接",  "銅色",    "連續金屬反光"],
    ]
    add_table(s, data, Inches(0.5), Inches(2.0), Inches(6.4), Inches(3.4),
              header_size=15, body_size=14,
              col_widths=[Inches(2.3), Inches(1.6), Inches(1.3), Inches(1.2)])

    # 基材說明（淺底簡單方塊）
    add_rect(s, Inches(0.5), Inches(5.7), Inches(6.4), Inches(1.3),
             fill=BG_ALT, line=GRAY_LIGHT)
    add_text(s, "基材特性",
             Inches(0.7), Inches(5.8), Inches(6.0), Inches(0.4),
             size=16, bold=True, color=NAVY)
    add_text(s, "FR-4 玻纖基板 ・ 綠色防焊漆 ・ 銀色焊墊 ・ 銅箔走線 ・ 黑色字符",
             Inches(0.7), Inches(6.25), Inches(6.1), Inches(0.8),
             size=14, color=DARK)

    # 右：缺陷對照圖
    s.shapes.add_picture(str(FIG / "fig1_defect_samples.png"),
                         Inches(7.2), Inches(1.6), height=Inches(5.4))

    add_footer(s, 2)


# ============================================================
# Slide 3 — 檢測方法
# ============================================================
def slide_3_method():
    s = blank_slide()
    add_title_bar(s, "檢測方法・系統架構")

    # 左：架構圖
    s.shapes.add_picture(str(FIG / "fig2_system_architecture.png"),
                         Inches(0.3), Inches(1.4), height=Inches(5.5))

    # 右：規格表
    add_text(s, "硬體規格", Inches(7.8), Inches(1.4),
             Inches(5.2), Inches(0.5), size=20, bold=True, color=NAVY)

    spec = [
        ["項目", "規格", "選用理由"],
        ["相機", "5MP CMOS", "0.3 mm 缺陷需 3 px"],
        ["鏡頭", "0.1× 遠心鏡頭", "消除 PCB 厚度透視變形"],
        ["主光", "同軸光", "金屬反射均勻明亮"],
        ["輔光", "低角度環形光", "邊緣陰影突顯缺陷"],
        ["觸發", "光遮斷感測器", "PCB 到位自動拍照"],
        ["機構", "暗箱 + 防震底座", "隔絕環境光、抑制振動"],
    ]
    add_table(s, spec, Inches(7.8), Inches(2.0), Inches(5.2), Inches(4.3),
              header_size=14, body_size=12,
              col_widths=[Inches(0.8), Inches(1.7), Inches(2.7)])

    # FOV 計算（簡潔灰底）
    add_rect(s, Inches(7.8), Inches(6.45), Inches(5.2), Inches(0.55),
             fill=BG_ALT, line=GRAY_LIGHT)
    add_text(s, "FOV 100×100 mm ・ 最小缺陷 0.3 mm ・ 解析度需求 ≥ 1 MP",
             Inches(7.95), Inches(6.5), Inches(5.0), Inches(0.45),
             size=13, bold=True, color=NAVY)

    add_footer(s, 3)


# ============================================================
# Slide 4 — 檢測流程
# ============================================================
def slide_4_flow():
    s = blank_slide()
    add_title_bar(s, "檢測流程")

    # 流程圖
    s.shapes.add_picture(str(FIG / "fig4_flow_chart.png"),
                         Inches(0.3), Inches(1.4), width=Inches(12.7))

    # 步驟說明（兩欄式）
    descriptions = [
        ("影像擷取", "光遮斷感測器觸發相機快門，同軸光與環形光同步亮起"),
        ("預處理", "原始影像縮放至 512×512 並正規化，符合模型輸入規格"),
        ("YOLOv12 推論", "GPU 單張延遲約 15 ms，可達 60+ FPS 即時檢測"),
        ("後處理", "Non-Maximum Suppression 去除重複框，依信心度過濾"),
        ("缺陷判定", "任一缺陷被偵測即判 FAIL，並記錄類別與位置"),
        ("輸出記錄", "結果寫入 log；連續 3 片 FAIL 觸發 SPC 警示"),
    ]
    for i, (head, body) in enumerate(descriptions):
        col = i % 2
        row = i // 2
        left = Inches(0.6 + col * 6.3)
        top = Inches(4.9 + row * 0.75)
        add_text(s, f"{i+1}. {head}",
                 left, top, Inches(2.2), Inches(0.4),
                 size=15, bold=True, color=NAVY)
        add_text(s, body,
                 left + Inches(2.2), top, Inches(4.0), Inches(0.45),
                 size=13, color=DARK)

    add_footer(s, 4)


# ============================================================
# Slide 5 — 影像處理
# ============================================================
def slide_5_image_proc():
    s = blank_slide()
    add_title_bar(s, "影像處理")

    # 上：擴增對照圖
    s.shapes.add_picture(str(FIG / "fig5_augmentation_demo.png"),
                         Inches(0.5), Inches(1.4), width=Inches(12.3))

    # 中：擴增策略說明（極簡橫向卡片）
    add_text(s, "資料擴增策略",
             Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5),
             size=22, bold=True, color=NAVY)
    add_text(s, "於軟體層補硬體未能涵蓋之變異",
             Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.45),
             size=15, color=GRAY)

    aug_items = [
        ("亮度 / 對比",  "模擬產線打光強度與焊墊反光變化"),
        ("水平翻轉",     "模擬 PCB 上料方向不一致"),
        ("旋轉 ±15°",    "模擬輸送帶到位角度誤差"),
    ]
    for i, (head, body) in enumerate(aug_items):
        left = Inches(0.5 + i * 4.3)
        add_rect(s, left, Inches(5.5), Inches(4.0), Inches(1.15),
                 fill=BG_ALT, line=GRAY_LIGHT)
        add_text(s, head, left + Inches(0.2), Inches(5.6),
                 Inches(3.8), Inches(0.45), size=17, bold=True, color=NAVY)
        add_text(s, body, left + Inches(0.2), Inches(6.05),
                 Inches(3.8), Inches(0.55), size=13, color=DARK)

    # 下：程式碼來源
    add_text(s,
             "程式碼來源：Ultralytics YOLOv12 ・ "
             "github.com/Joshua-fy-Hsu/PCB_Defects_Detection_Using_YOLOv12",
             Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
             size=11, color=GRAY)

    add_footer(s, 5)


# ============================================================
# Slide 6 — YOLOv12 模型介紹
# ============================================================
def slide_6_yolov12():
    s = blank_slide()
    add_title_bar(s, "YOLOv12 模型介紹")

    add_text(s, "本專題採用之深度學習偵測模型 — 2025 年最新 SOTA 即時物件偵測架構",
             Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             size=18, color=GRAY)

    # ---- 左：模型特色 ----
    add_text(s, "模型特色",
             Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.5),
             size=22, bold=True, color=NAVY)
    features = [
        "Attention-Centric 架構：首次將自注意力機制\n      整合進 YOLO 系列",
        "Area Attention (AAttn)：低計算成本之區域\n      注意力，比傳統 self-attention 快數倍",
        "R-ELAN 骨幹：殘差結構強化梯度傳遞，提升\n      訓練穩定性與小物件偵測能力",
        "Anchor-free One-stage Detector：免設計 anchor\n      box，預測更直接",
        "在 COCO 公開資料集上同時取得最佳 mAP\n      與 FPS 表現",
    ]
    add_bullets(s, features, Inches(0.5), Inches(2.55),
                Inches(6.5), Inches(4.3), size=14, space_before=10)

    # ---- 右：訓練配置 ----
    add_text(s, "訓練配置",
             Inches(7.3), Inches(2.0), Inches(5.7), Inches(0.5),
             size=22, bold=True, color=NAVY)
    cfg = [
        ["項目", "設定"],
        ["模型版本", "YOLOv12s (small)"],
        ["預訓練權重", "yolov12s.pt (COCO)"],
        ["輸入尺寸", "512 × 512"],
        ["訓練 Epochs", "100"],
        ["Batch Size", "16"],
        ["最佳化演算法", "AdamW"],
        ["訓練裝置", "NVIDIA RTX 4060 (8GB)"],
    ]
    add_table(s, cfg, Inches(7.3), Inches(2.55),
              Inches(5.7), Inches(4.0),
              header_size=14, body_size=13,
              col_widths=[Inches(2.0), Inches(3.7)])

    # ---- 底部說明條 ----
    add_rect(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.55),
             fill=BG_ALT, line=GRAY_LIGHT)
    add_text(s,
             "選用理由：相較傳統 CNN-only 的 YOLO 系列，"
             "注意力機制有助於在複雜 PCB 背景中聚焦於缺陷區域，提高小目標召回率。",
             Inches(0.7), Inches(6.75), Inches(12.0), Inches(0.45),
             size=13, color=NAVY)

    add_footer(s, 6)


# ============================================================
# Slide 7 — 檢測結果
# ============================================================
def slide_7_results():
    s = blank_slide()
    add_title_bar(s, "檢測結果")

    # 左：指標長條圖
    s.shapes.add_picture(str(FIG / "fig6_metrics_bar.png"),
                         Inches(0.3), Inches(1.4), width=Inches(7.0))

    # 右上：混淆矩陣
    cm_path = REPO / "outputs" / "enhanced" / "confusion_matrix.png"
    if cm_path.exists():
        s.shapes.add_picture(str(cm_path), Inches(7.6), Inches(1.4),
                             height=Inches(3.5))
    add_text(s, "混淆矩陣  Confusion Matrix",
             Inches(7.6), Inches(4.95), Inches(5.3), Inches(0.4),
             size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # 下：dashboard demo 區塊（極簡）
    add_rect(s, Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.55),
             fill=BG_ALT, line=NAVY, line_w=1.5)
    add_text(s, "模擬產線 Dashboard 即時展示",
             Inches(0.6), Inches(5.7), Inches(12.3), Inches(0.5),
             size=20, bold=True, color=NAVY)
    add_text(s,
             "以 Streamlit 模擬 20 片 PCB 依序進入檢測站，即時顯示 PASS/FAIL、"
             "累積良率趨勢、缺陷分布與連續不良 SPC 警示，並自動寫入 CSV。",
             Inches(0.6), Inches(6.2), Inches(12.3), Inches(0.55),
             size=14, color=DARK)
    add_text(s, "現場 demo：streamlit run code/dashboard.py",
             Inches(0.6), Inches(6.75), Inches(12.3), Inches(0.4),
             size=13, bold=True, color=NAVY)

    add_footer(s, 7)


# ============================================================
# Slide 8 — 專題應用（簡潔卡片、無 icon）
# ============================================================
def slide_8_applications():
    s = blank_slide()
    add_title_bar(s, "專題應用")
    add_text(s, "本套 AOI 流程可延伸至其他製造業缺陷檢測場域",
             Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             size=20, color=GRAY)

    apps = [
        ("SMT 焊後檢測",     "錫珠、空焊、零件偏移、立碑"),
        ("IC 封裝外觀",      "崩角、刮傷、字符 OCR 比對"),
        ("觸控面板 ITO 線路", "細線路斷路 / 短路檢測"),
        ("金屬沖壓件表面",   "毛邊、凹痕、刮傷"),
        ("太陽能板隱裂",     "EL 影像中細微裂紋偵測"),
        ("玻璃 / 光學鏡片",  "氣泡、刮痕、汙點"),
    ]
    card_w, card_h = Inches(4.0), Inches(2.3)
    gap = Inches(0.18)
    start_x = Inches(0.55)
    start_y = Inches(2.05)

    for i, (title, desc) in enumerate(apps):
        col = i % 3
        row = i // 3
        left = start_x + col * (card_w + gap)
        top = start_y + row * (card_h + gap)
        add_rect(s, left, top, card_w, card_h, fill=WHITE, line=NAVY, line_w=1.5)
        add_text(s, title, left + Inches(0.25), top + Inches(0.25),
                 card_w - Inches(0.5), Inches(0.6),
                 size=20, bold=True, color=NAVY)
        # 細分隔線
        sub = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                 left + Inches(0.25), top + Inches(0.95),
                                 Inches(0.7), Inches(0.03))
        sub.fill.solid(); sub.fill.fore_color.rgb = NAVY
        sub.line.fill.background(); sub.shadow.inherit = False
        add_text(s, desc, left + Inches(0.25), top + Inches(1.15),
                 card_w - Inches(0.5), Inches(1.0),
                 size=15, color=DARK)

    add_footer(s, 8)


# ============================================================
# Slide 9 — 與課程相關知識
# ============================================================
def slide_9_course():
    s = blank_slide()
    add_title_bar(s, "專題製作與課程相關知識")

    # 左：打光對比圖
    s.shapes.add_picture(str(FIG / "fig3_lighting_comparison.png"),
                         Inches(0.3), Inches(1.5), width=Inches(6.4))
    add_text(s, "打光方式對應課程「工業照明」章節",
             Inches(0.3), Inches(4.6), Inches(6.4), Inches(0.4),
             size=13, color=GRAY, align=PP_ALIGN.CENTER)

    # 右：對應表
    add_text(s, "課程主題 對應", Inches(7.0), Inches(1.5),
             Inches(6.0), Inches(0.5), size=20, bold=True, color=NAVY)
    course = [
        ["課程主題", "本專題對應"],
        ["鏡頭選型與 FOV 計算", "遠心鏡頭 + 5 MP 解析度推算"],
        ["工業打光技術", "同軸光 + 低角度環形光"],
        ["影像前處理", "Resize / Normalize / Augmentation"],
        ["特徵擷取與分類", "YOLOv12 卷積骨幹網路"],
        ["缺陷檢測演算法", "One-stage Anchor-free Detector"],
        ["系統整合與部署", "觸發、流程、PASS/FAIL、SPC 警示"],
    ]
    add_table(s, course, Inches(7.0), Inches(2.05),
              Inches(6.0), Inches(4.4),
              header_size=15, body_size=13,
              col_widths=[Inches(2.4), Inches(3.6)])

    # 底部摘要
    add_rect(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.55),
             fill=BG_ALT, line=GRAY_LIGHT)
    add_text(s,
             "從「鏡頭與打光的硬體設計」到「演算法與系統整合」，"
             "完整涵蓋本課程機器視覺檢測知識體系。",
             Inches(0.7), Inches(6.65), Inches(12.0), Inches(0.45),
             size=14, color=NAVY)

    add_footer(s, 9)


# ============================================================
# Slide 10 — 心得 & 貢獻（許豐有為主要負責人）
# ============================================================
def slide_10_reflection():
    s = blank_slide()
    add_title_bar(s, "專題製作心得 & 貢獻")

    members = [
        {
            "id": "411287004",
            "name": "許豐有",
            "role": "（主要負責人）",
            "contrib": [
                "資料蒐集、標註與前處理",
                "模型訓練與超參數調校",
                "系統架構與打光設計",
                "Streamlit Dashboard 開發",
                "效能評估與結果分析",
                "簡報撰寫與報告整合",
            ],
            "reflect": (
                "親自走完從資料、訓練到部署模擬的整條流程後，最深的體會是"
                "「軟體層的資料擴增」其實是在補硬體未能涵蓋的變異；"
                "產線檢測的可靠度來自軟硬體共同設計，而非單靠模型本身。"
            ),
        },
        {
            "id": "411287035",
            "name": "趙梓辰",
            "role": "（協作）",
            "contrib": [
                "資料標註輔助",
                "簡報部分章節協作",
            ],
            "reflect": (
                "透過協助專題理解到 PCB 缺陷的多樣性，以及在現代產線中導入"
                "機器視覺對品質一致性的重要性。"
            ),
        },
    ]

    card_w = Inches(6.1)
    for i, m in enumerate(members):
        left = Inches(0.4 + i * 6.45)
        # 卡片底框
        add_rect(s, left, Inches(1.45), card_w, Inches(5.55),
                 fill=WHITE, line=NAVY, line_w=1.5)
        # 頂部 navy 條
        head = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                  left, Inches(1.45), card_w, Inches(0.85))
        head.fill.solid(); head.fill.fore_color.rgb = NAVY
        head.line.fill.background(); head.shadow.inherit = False
        add_text(s, f"{m['id']}   {m['name']}",
                 left + Inches(0.3), Inches(1.55),
                 card_w - Inches(0.6), Inches(0.5),
                 size=20, bold=True, color=WHITE)
        add_text(s, m["role"],
                 left + Inches(0.3), Inches(1.95),
                 card_w - Inches(0.6), Inches(0.35),
                 size=12, color=GRAY_LIGHT)

        # 貢獻
        add_text(s, "製作貢獻",
                 left + Inches(0.3), Inches(2.5),
                 card_w - Inches(0.6), Inches(0.4),
                 size=16, bold=True, color=NAVY)
        bullet_h = 0.45 * len(m["contrib"]) + 0.3
        add_bullets(s, m["contrib"],
                    left + Inches(0.3), Inches(2.95),
                    card_w - Inches(0.6), Inches(bullet_h),
                    size=14, space_before=6)

        # 心得
        reflect_top = Inches(2.95 + bullet_h)
        add_text(s, "製作心得",
                 left + Inches(0.3), reflect_top,
                 card_w - Inches(0.6), Inches(0.4),
                 size=16, bold=True, color=NAVY)
        add_text(s, m["reflect"],
                 left + Inches(0.3), reflect_top + Inches(0.45),
                 card_w - Inches(0.6), Inches(2.0),
                 size=13, color=DARK)

    add_footer(s, 10)


# ============================================================
# Slide 11 — 參考資料
# ============================================================
def slide_11_references():
    s = blank_slide()
    add_title_bar(s, "參考資料")

    refs = [
        ("[1] Ultralytics YOLOv12 (Sun et al., GitHub)",
         "https://github.com/sunsmarterjie/yolov12"),
        ("[2] Ultralytics 官方文件",
         "https://docs.ultralytics.com"),
        ("[3] PCB Defects Dataset — Kaggle (Akhatova)",
         "https://www.kaggle.com/datasets/akhatova/pcb-defects"),
        ("[4] Roboflow PCB Defects — pcb-nzihq workspace",
         "https://universe.roboflow.com/pcb-nzihq/pcb-defects-kmf3o"),
        ("[5] Streamlit — Open-source Python App Framework",
         "https://streamlit.io"),
        ("[6] Basler 工業相機產品線",
         "https://www.baslerweb.com"),
        ("[7] Edmund Optics — Telecentric Lens 應用筆記",
         "https://www.edmundoptics.com"),
        ("[8] 本專題 GitHub Repository",
         "https://github.com/Joshua-fy-Hsu/PCB_Defects_Detection_Using_YOLOv12"),
        ("[9] 機器視覺檢測課程講義（鏡頭、打光、影像處理章節）",
         "—"),
    ]

    for i, (title, url) in enumerate(refs):
        top = Inches(1.45 + i * 0.6)
        add_text(s, title, Inches(0.6), top, Inches(7.5), Inches(0.4),
                 size=15, bold=True, color=DARK)
        add_text(s, url, Inches(0.6), top + Inches(0.32),
                 Inches(12.0), Inches(0.3), size=12, color=GRAY)

    add_footer(s, 11)


# ============================================================
# 產出
# ============================================================
if __name__ == "__main__":
    slide_1_cover()
    slide_2_object()
    slide_3_method()
    slide_4_flow()
    slide_5_image_proc()
    slide_6_yolov12()
    slide_7_results()
    slide_8_applications()
    slide_9_course()
    slide_10_reflection()
    slide_11_references()

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"[OK] {OUT}")
